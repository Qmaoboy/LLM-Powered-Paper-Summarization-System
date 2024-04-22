from pptx import Presentation
from pptx.util import Inches,Pt
from pptx.enum.text import PP_ALIGN,MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from datetime import datetime
import re,os
from utilized import *
from PIL import Image
from io import BytesIO

def make_ppt(Target_Folder,Paper_list,client,args):

    ppt_path=os.path.join(Target_Folder.strip(),'PPT')
    os.makedirs(ppt_path,exist_ok=True)

    ## Watermark
    def add_watermark(slide):
        watermark_text = "Nycu MLlab"
        watermark_font_size = 20
        watermark_color = RGBColor(200, 200, 200)  # 淺灰色
        left,top,width,height=Inches(8), Inches(6.5), Inches(2), Inches(0.5)
        text_box = slide.shapes.add_textbox(left,top,width,height)
        rotation_angle = -45
        text_box.rotation = rotation_angle
        text_frame = text_box.text_frame
        p = text_frame.paragraphs[0]
        p.text = watermark_text
        p.font.size = Pt(watermark_font_size)
        p.font.color.rgb = watermark_color
        p.alignment = PP_ALIGN.CENTER

        # 調整文字框的透明度
        # text_box.fill.solid()
        # text_box.fill.fore_color.rgb = RGBColor(255, 255, 255)  # 白色背景
        text_box.fill.transparency = 1.0  # 透明度 0.5 (0 = 不透明, 1 = 完全透明)
        # 設定文字方塊的邊框為透明

        ## add Logo
        left,top,width,height= Inches(0.1),Inches(6.5),Inches(1),Inches(1)
        img_path=r"backend\app\assets\InnoLux_Corporation-Logo.wine.png"
        pic = slide.shapes.add_picture(img_path, left,top,width,height)
        return slide

    papertitle=re.sub(r"\s+", "", ppt_path.split('\\')[-2].split(".")[0])
    version_name="3.8"
    key=""
    X=Presentation()
    Layout = X.slide_layouts[0]
    first_slide = X.slides.add_slide(Layout) # Adding first slide
    first_slide=add_watermark(first_slide)
    first_slide.shapes.title.text = "Innolux Paper Explaination \nwith ChatGPT"
    first_slide.placeholders[1].text = "NYCU MLLAB\n"+papertitle

    txBox0 = first_slide.shapes.add_textbox(Inches(9),Inches(7),Inches(1),Inches(0.3))
    inittf = txBox0.text_frame
    inittf.paragraphs[0].text=f"Version: {version_name}"
    inittf.paragraphs[0].alignment = PP_ALIGN.RIGHT
    inittf.vertical_anchor = MSO_ANCHOR.TOP
    inittf.word_wrap=True
     ## Title font
    font = inittf.paragraphs[0].font
    # font.name = '微软雅黑'    # 字体类型
    font.name = 'Verdana'
    font.bold = True    # 加粗
    font.size = Pt(10)    # 大小
    paper_layout=X.slide_layouts[6]
    for key in Paper_list:
        # try:
        value=client.find_data("Papers_info","Paper_name",key)
        if not value:
            # print(f"{key} Not Found in Mysql Papers_info Table")
            continue
        value=value.pop()
        slide_ = X.slides.add_slide(paper_layout)
        slide_=add_watermark(slide_)
        ## Title textbox

        txBox0 = slide_.shapes.add_textbox(Inches(0.1),Inches(0.1),Inches(9.9),Inches(1))
        tf = txBox0.text_frame

        tf.paragraphs[0].text=value['title_name']
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.add_paragraph()
        tf.paragraphs[1].text=value['author_name']+'\n'+value['organization_name']
        tf.paragraphs[1].alignment = PP_ALIGN.CENTER
        ## Title font
        font = tf.paragraphs[0].font
        # font.name = '微软雅黑'    # 字体类型
        font.bold = True    # 加粗
        font.size = Pt(15)    # 大小
        tf.word_wrap=True

        font = tf.paragraphs[1].font
        # font.name = '微软雅黑'    # 字体类型
        font.bold = True    # 加粗
        font.size = Pt(10)    # 大小
        tf.word_wrap=True

        ## keyword textbox
        left,top,width,height=Inches(0.1),Inches(1.2),Inches(9.9),Inches(0.5)
        txBox2 = slide_.shapes.add_textbox(left, top, width, height)
        tf2 = txBox2.text_frame
        tf2.word_wrap=True
        tf2.paragraphs[0].text = "Keywords : "+value['keywords']+"\nOrganization: "+value['organization_name']
        font1 = tf2.paragraphs[0].font
        # font.name = '微软雅黑'    # 字体类型
        font1.bold = True    # 加粗
        font1.size = Pt(10)    # 大小

        ## GPT reply textbox
        left,top,width,height=Inches(0.1),Inches(2),Inches(9.9),Inches(2)
        txBox1 = slide_.shapes.add_textbox(left, top, width, height)
        tf1 = txBox1.text_frame
        tf1.word_wrap=True

        tf1.paragraphs[0].text = f"Keypoints : \n- "+value['keypoints'].replace(",", "\n- ")
        ## First textbox font
        font0 = tf1.paragraphs[0].font
        # font.name = '微软雅黑'    # 字体类型
        font0.bold = True    # 加粗
        font0.size = Pt(11)    # 大小

        ## Setup How many Image show on PPT
        pagcount=3
        # print(f'{value["Paper_name"]}: {image_count} pictures')
        if value['Represent_image'] is None:
            print(f"{key} Represent_image Not Found in Mysql Papers_info Table")
            continue
        ppt_image=value['Represent_image']
        ppt_image_count=len(ppt_image)
        if ppt_image_count:
            if ppt_image_count>=pagcount:
                ppt_image_count=pagcount
            for i in range(ppt_image_count):
                ll=(10-0.1*(ppt_image_count+1))/ppt_image_count
                img_name=ppt_image[i]
                imgid=client.find_data("Image_Files","Image_id",img_name)
                if imgid:
                    imgid=imgid.pop()["Image_Byte_File"]
                    image = Image.open(BytesIO(imgid))## 開起圖片
                    img_bytesio = BytesIO()
                    image.save(img_bytesio, format='PNG')
                    left,top,width,height=Inches(0.1*(i+1)+ll*i),Inches(4.2),Inches(ll),Inches(2)
                    slide_.shapes.add_picture(img_bytesio,left,top,width,height)
                    ## BLIp prompt
                    left,top,width,height=Inches(0.1*(i+1)+ll*i),Inches(6.2),Inches(ll),Inches(1)
                    txBoxi = slide_.shapes.add_textbox(left, top, width, height)
                    tfi = txBoxi.text_frame
                    tfi.word_wrap=True
                    if img_name in value['image_Description']:
                        tfi.paragraphs[0].text =value['image_Description'][img_name]
                        fonti = tfi.paragraphs[0].font
                        # font.name = '微软雅黑'    # 字体类型
                        fonti.bold = True    # 加粗
                        fonti.size = Pt(11)    # 大小
                else:
                    print(f"{img_name} Not Found in Mysql Image_Files Table")
    now = datetime.now()
    now_time = now.strftime("%Y_%m_%d_%H_%M")
    target_file=f"{ppt_path}\Innolux_paper_{papertitle}_{now_time}.pptx"
    X.save(target_file)

    return target_file
